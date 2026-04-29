"""Builds a .pptx territory report from slide records and S3 images."""

from io import BytesIO
from collections.abc import Generator, Sequence
from typing import Any

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # type: ignore[attr-defined]
from pptx.util import Inches, Pt

from src.models.exports import MapExportSlideModel
from src.services.s3 import S3Service

# Slide dimensions — standard 16:9
_W = Inches(10)
_H = Inches(5.625)

# Layout constants
_TITLE_H = Inches(0.55)
_CONTENT_Y = _TITLE_H
_CONTENT_H = _H - _TITLE_H
_IMAGE_W = Inches(10 * 2 / 3)  # left 2/3
_TABLE_X = _IMAGE_W
_TABLE_W = _W - _IMAGE_W       # right 1/3

_HEADER_COLOR = RGBColor(0xE5, 0xE7, 0xEB)  # gray-200


def _fit_in_box(img_bytes: bytes, box_w: int, box_h: int) -> tuple[int, int]:
    """Return (width, height) in EMUs that fit within box_w x box_h preserving aspect ratio."""
    with PILImage.open(BytesIO(img_bytes)) as img:
        iw, ih = img.size
    scale = min(box_w / iw, box_h / ih)
    return round(iw * scale), round(ih * scale)


def _add_slide(
    prs: Presentation,
    blank_layout: Any,
    title: str,
    image_bytes: bytes | None,
    node_data: list[dict[str, Any]],
) -> None:
    slide = prs.slides.add_slide(blank_layout)

    # Title
    tb = slide.shapes.add_textbox(0, 0, _W, _TITLE_H)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.LEFT

    # Map screenshot — fit within the slot preserving the captured aspect ratio
    if image_bytes:
        img_w, img_h = _fit_in_box(image_bytes, int(_IMAGE_W), int(_CONTENT_H))
        slide.shapes.add_picture(
            BytesIO(image_bytes),
            left=0,
            top=_CONTENT_Y,
            width=img_w,
            height=img_h,
        )

    # Table
    if not node_data:
        return

    first = node_data[0]
    data_keys = [k for k in first if k != "name"]
    raw_keys = ["name"] + data_keys
    col_headers = ["Name"] + [k.replace("_", " ").title() for k in data_keys]

    tbl = slide.shapes.add_table(
        len(node_data) + 1,
        len(col_headers),
        _TABLE_X,
        _CONTENT_Y,
        _TABLE_W,
        _CONTENT_H,
    ).table

    # Header row
    for ci, header in enumerate(col_headers):
        cell = tbl.cell(0, ci)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = _HEADER_COLOR
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(9)
        para.alignment = PP_ALIGN.CENTER

    # Data rows
    for ri, node in enumerate(node_data, start=1):
        for ci, key in enumerate(raw_keys):
            val = node.get(key)
            cell = tbl.cell(ri, ci)
            cell.text = "" if val is None else str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(8)


def build_pptx_chunks(
    slides: Sequence[MapExportSlideModel],
    s3: S3Service,
    chunk_size: int = 65536,
) -> Generator[bytes, None, None]:
    """Build the .pptx in memory (one S3 image at a time) then yield it in chunks."""
    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H
    blank_layout = prs.slide_layouts[6]

    for slide_model in slides:
        image_bytes: bytes | None = None
        if slide_model.image_s3_key:
            image_bytes = s3.get_private_object(key=slide_model.image_s3_key).read()

        _add_slide(
            prs,
            blank_layout,
            title=slide_model.title or "",
            image_bytes=image_bytes,
            node_data=slide_model.node_data or [],
        )

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    while chunk := buf.read(chunk_size):
        yield chunk
